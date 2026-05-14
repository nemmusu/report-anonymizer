"""PPTX adapter built on python-pptx.

One segment per ``Run`` of every ``TextFrame`` (slide shapes, table cells,
notes_slide). Same run-merge strategy as :class:`DocxAdapter`.
"""
from __future__ import annotations

from pathlib import Path
from typing import Iterable

from .base import (
    FormatAdapter,
    ImageReport,
    InventoryImageRaw,
    Segment,
    SubstitutionRule,
    WriteEvent,
    WriteReport,
    apply_to_text,
)


def _iter_text_frames(prs) -> Iterable[tuple[str, object]]:
    for si, slide in enumerate(prs.slides):
        prefix = f"slide{si}"
        yield from _iter_in_shapes(prefix, slide.shapes)
        if slide.has_notes_slide:
            ns = slide.notes_slide
            if ns is not None:
                yield from _iter_in_shapes(f"{prefix}.notes", ns.shapes)


def _iter_in_shapes(prefix: str, shapes) -> Iterable[tuple[str, object]]:
    for shi, shape in enumerate(shapes):
        if shape.has_text_frame:
            yield f"{prefix}.sh{shi}.tf", shape.text_frame
        if shape.shape_type == 19 or getattr(shape, "has_table", False):
            try:
                tbl = shape.table
            except Exception:
                continue
            for ri, row in enumerate(tbl.rows):
                for ci, cell in enumerate(row.cells):
                    if cell.text_frame is not None:
                        yield f"{prefix}.sh{shi}.tbl.r{ri}c{ci}", cell.text_frame
        if hasattr(shape, "shapes"):
            try:
                yield from _iter_in_shapes(f"{prefix}.sh{shi}.grp", shape.shapes)
            except Exception:
                pass


def _redistribute_runs(runs: list, full_text: str) -> None:
    pos = 0
    n = len(full_text)
    for idx, run in enumerate(runs):
        original_len = len(run.text or "")
        if pos >= n:
            run.text = ""
            continue
        if idx == len(runs) - 1:
            run.text = full_text[pos:]
            return
        if original_len == 0:
            continue
        end = min(pos + original_len, n)
        run.text = full_text[pos:end]
        pos = end


class PptxAdapter(FormatAdapter):
    name = "pptx"
    extensions = {".pptx", ".pptm"}
    mimes = {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }

    def __init__(self) -> None:
        try:
            from pptx import Presentation  # type: ignore
        except Exception as e:  # pragma: no cover
            raise RuntimeError(
                f"python-pptx is required for the pptx adapter: {e}"
            )
        self._Presentation = Presentation

    def extract(self, path: Path) -> list[Segment]:
        prs = self._Presentation(str(path))
        out: list[Segment] = []
        for tf_id, tf in _iter_text_frames(prs):
            for pi, para in enumerate(tf.paragraphs):
                text = "".join(r.text or "" for r in para.runs)
                if not text:
                    continue
                out.append(Segment(seg_id=f"{tf_id}.p{pi}", text=text))
        return out

    def write(
        self,
        src_path: Path,
        dst_path: Path,
        substitutions: list[SubstitutionRule],
    ) -> WriteReport:
        prs = self._Presentation(str(src_path))
        events: list[WriteEvent] = []
        for tf_id, tf in _iter_text_frames(prs):
            for pi, para in enumerate(tf.paragraphs):
                runs = para.runs
                if not runs:
                    continue
                full = "".join(r.text or "" for r in runs)
                if not full:
                    continue
                new, ev = apply_to_text(
                    full, substitutions, seg_id=f"{tf_id}.p{pi}"
                )
                if new == full:
                    continue
                _redistribute_runs(runs, new)
                events.extend(ev)
        dst_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(dst_path))
        return WriteReport(file_rel=str(dst_path), events=events)

    # ---- image inventory + apply ------------------------------------------

    def inventory_images(self, path: Path) -> list[InventoryImageRaw]:
        """Walk every slide, collect picture shapes + image parts.

        Pictures share the underlying ``ImagePart`` when the same
        image is reused across slides (the common case, e.g. a logo
        on every slide), so we inventory by ``(slide_index, shape_id)``
        for the GUI gallery while remembering the rel id for apply.
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
        except Exception:
            MSO_SHAPE_TYPE = None  # type: ignore
        try:
            from PIL import Image as _PIL  # type: ignore
        except Exception:
            _PIL = None  # type: ignore

        out: list[InventoryImageRaw] = []
        try:
            prs = self._Presentation(str(path))
        except Exception:
            return out

        seen_part_keys: set[int] = set()

        for si, slide in enumerate(prs.slides):
            for shape in _iter_picture_shapes(slide.shapes, MSO_SHAPE_TYPE):
                try:
                    image = shape.image
                except Exception:
                    continue
                try:
                    raw = bytes(image.blob)
                except Exception:
                    continue
                # python-pptx exposes the underlying part via image.ext
                # only indirectly; use ``id(image_part)`` as the
                # process-local dedup key for raw_bytes when reading.
                part_id = id(getattr(image, "_blob", image))
                if part_id in seen_part_keys:
                    continue
                seen_part_keys.add(part_id)

                fmt = (image.ext or "png").lower()
                width, height = _pptx_image_size(image, _PIL)
                rel_id = _pptx_picture_rel_id(shape)
                out.append(
                    InventoryImageRaw(
                        raw_bytes=raw,
                        fmt=fmt,
                        width=width,
                        height=height,
                        location={
                            "kind": "pptx",
                            "slide_index": si,
                            "shape_id": int(shape.shape_id),
                            "rel_id": rel_id,
                        },
                    )
                )
        return out

    def apply_image_redactions(
        self,
        dst_path: Path,
        decisions_for_file: dict,
    ) -> ImageReport:
        """Replace ``image_part.blob`` for every picture whose
        ``image_id`` has a ``redact`` decision. Saves atomically via
        a temp file in the same directory plus ``os.replace``.
        """
        from ..image_inventory import compute_image_id, ImageDecision
        from ..image_redactor import ImageRedaction, ImageRedactor
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
        except Exception:
            MSO_SHAPE_TYPE = None  # type: ignore

        report = ImageReport(file_rel=str(dst_path))
        if not decisions_for_file:
            return report

        prepared: dict[str, list] = {}
        skip_ids: set[str] = set()
        for image_id, decision in decisions_for_file.items():
            if not isinstance(decision, ImageDecision):
                decision = ImageDecision.from_dict(image_id, decision or {})
            if decision.decision in ("skip", "defer"):
                skip_ids.add(image_id)
                continue
            if decision.decision != "redact" or not decision.rects:
                skip_ids.add(image_id)
                continue
            prepared[image_id] = [
                ImageRedaction(
                    x=r.x, y=r.y, w=r.w, h=r.h,
                    tool=r.tool,
                    intensity=r.intensity,
                    text=r.text,
                    font_size=r.font_size,
                    fg=r.fg,
                    bg=r.bg,
                ) for r in decision.rects
            ]
        if not prepared and not skip_ids:
            return report

        try:
            prs = self._Presentation(str(dst_path))
        except Exception as e:
            report.warnings.append(f"open_failed:{e}")
            return report

        seen_part_objs: set[int] = set()
        for slide in prs.slides:
            for shape in _iter_picture_shapes(slide.shapes, MSO_SHAPE_TYPE):
                try:
                    image = shape.image
                except Exception:
                    continue
                # Resolve to the underlying ImagePart so we can mutate
                # its blob; setting shape.image.blob does not exist as
                # a writable attribute in python-pptx.
                image_part = _pptx_image_part(slide, shape)
                if image_part is None:
                    continue
                if id(image_part) in seen_part_objs:
                    continue
                seen_part_objs.add(id(image_part))
                try:
                    raw = bytes(image.blob)
                except Exception:
                    continue
                image_id = compute_image_id(raw)
                if image_id in skip_ids:
                    report.skipped += 1
                    continue
                rects = prepared.get(image_id)
                if not rects:
                    report.untouched += 1
                    continue
                fmt_hint = (image.ext or "png").lower()
                try:
                    result = ImageRedactor.redact_bytes(raw, fmt_hint, rects)
                except Exception as e:
                    report.warnings.append(f"redact_failed:{image_id}:{e}")
                    continue
                try:
                    image_part._blob = result.bytes_
                except Exception as e:
                    report.warnings.append(f"replace_failed:{image_id}:{e}")
                    continue
                if result.warnings:
                    report.warnings.extend(
                        f"{image_id}:{w}" for w in result.warnings
                    )
                report.applied += 1

        tmp_save = dst_path.with_suffix(dst_path.suffix + ".imgtmp")
        try:
            prs.save(str(tmp_save))
        except Exception as e:
            report.warnings.append(f"save_failed:{e}")
            try:
                tmp_save.unlink()
            except OSError:
                pass
            return report
        import os as _os
        _os.replace(str(tmp_save), str(dst_path))
        return report


def _iter_picture_shapes(shapes, MSO_SHAPE_TYPE):
    """Yield every shape that wraps an embedded picture.

    Walks group shapes recursively: pptx allows pictures nested
    inside groups, and ``slide.shapes`` does NOT expand them by
    default.
    """
    for shape in shapes:
        try:
            stype = shape.shape_type
        except Exception:
            stype = None
        if MSO_SHAPE_TYPE is not None and stype == MSO_SHAPE_TYPE.PICTURE:
            yield shape
        elif MSO_SHAPE_TYPE is not None and stype == MSO_SHAPE_TYPE.GROUP:
            try:
                yield from _iter_picture_shapes(shape.shapes, MSO_SHAPE_TYPE)
            except Exception:
                continue
        elif hasattr(shape, "image"):
            # Belt + braces: some pptx variants tag pictures with
            # other shape types. If the shape exposes ``image``, we
            # treat it as a picture.
            yield shape


def _pptx_picture_rel_id(shape) -> str:
    """Return the relationship id of the image referenced by ``shape``."""
    try:
        return str(shape._element.blip_rId)
    except Exception:
        try:
            return str(shape._pic.blip_rId)
        except Exception:
            return ""


def _pptx_image_part(slide, shape):
    """Resolve the ImagePart behind a picture shape on a given slide."""
    rel_id = _pptx_picture_rel_id(shape)
    if not rel_id:
        return None
    try:
        rel = slide.part.rels[rel_id]
    except Exception:
        return None
    return getattr(rel, "target_part", None)


def _pptx_image_size(image, pil_module) -> tuple[int, int]:
    if pil_module is None:
        return (0, 0)
    try:
        from io import BytesIO
        img = pil_module.open(BytesIO(image.blob))
        img.load()
        return (int(img.size[0]), int(img.size[1]))
    except Exception:
        return (0, 0)


__all__ = ["PptxAdapter"]
