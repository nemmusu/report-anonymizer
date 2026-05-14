"""Round-trip integration tests for the PPTX image-redaction path."""
from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")
PIL = pytest.importorskip("PIL")
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from anonymize.format_adapters.pptx_adapter import PptxAdapter
from anonymize.image_inventory import (
    ImageDecision,
    RedactionRect,
    compute_image_id,
)


def _png_solid(color: tuple[int, int, int], size: tuple[int, int] = (200, 200)) -> bytes:
    img = Image.new("RGB", size, color)
    buf = BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_pptx(path: Path, image_bytes: bytes, *, slide_count: int = 1) -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for _ in range(slide_count):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(BytesIO(image_bytes), Inches(1), Inches(1),
                                 Inches(2), Inches(2))
    prs.save(str(path))


def test_inventory_lists_one_image_with_metadata(tmp_path: Path) -> None:
    src = _png_solid((220, 30, 30))
    p = tmp_path / "deck.pptx"
    _make_pptx(p, src)

    inv = PptxAdapter().inventory_images(p)
    assert len(inv) == 1
    item = inv[0]
    assert item.location["kind"] == "pptx"
    assert item.location["slide_index"] == 0
    assert item.location["shape_id"] is not None
    assert item.fmt == "png"
    assert item.width == 200 and item.height == 200


def test_skip_decision_keeps_blob_byte_identical(tmp_path: Path) -> None:
    src = _png_solid((30, 200, 30))
    p = tmp_path / "skip.pptx"
    _make_pptx(p, src)

    adapter = PptxAdapter()
    inv = adapter.inventory_images(p)
    image_id = compute_image_id(inv[0].raw_bytes)

    report = adapter.apply_image_redactions(
        p, {image_id: ImageDecision(image_id=image_id, decision="skip")}
    )
    assert report.applied == 0
    assert report.skipped == 1
    inv2 = PptxAdapter().inventory_images(p)
    assert compute_image_id(inv2[0].raw_bytes) == image_id


def test_redact_blackout_changes_pixels_inside_rect_only(tmp_path: Path) -> None:
    src = _png_solid((220, 30, 30))
    p = tmp_path / "redact.pptx"
    _make_pptx(p, src)

    adapter = PptxAdapter()
    inv = adapter.inventory_images(p)
    image_id = compute_image_id(inv[0].raw_bytes)

    report = adapter.apply_image_redactions(p, {
        image_id: ImageDecision(
            image_id=image_id,
            decision="redact",
            rects=[RedactionRect(x=50, y=50, w=100, h=100, tool="blackout")],
        )
    })
    assert report.applied == 1

    inv2 = PptxAdapter().inventory_images(p)
    new_img = Image.open(BytesIO(inv2[0].raw_bytes))
    new_img.load()
    assert new_img.size == (200, 200)
    assert new_img.getpixel((100, 100)) == (0, 0, 0)
    assert new_img.getpixel((10, 10)) == (220, 30, 30)


def test_logo_on_three_slides_redacted_once(tmp_path: Path) -> None:
    src = _png_solid((30, 30, 200))
    p = tmp_path / "logo.pptx"
    _make_pptx(p, src, slide_count=3)

    adapter = PptxAdapter()
    inv = adapter.inventory_images(p)
    # Same bytes on all 3 slides -> single image_id.
    ids = {compute_image_id(im.raw_bytes) for im in inv}
    assert len(ids) == 1
    image_id = next(iter(ids))

    adapter.apply_image_redactions(p, {
        image_id: ImageDecision(
            image_id=image_id,
            decision="redact",
            rects=[RedactionRect(x=0, y=0, w=200, h=50, tool="blackout")],
        )
    })

    # Reopen and assert every occurrence is redacted.
    inv2 = PptxAdapter().inventory_images(p)
    for item in inv2:
        img = Image.open(BytesIO(item.raw_bytes))
        img.load()
        assert img.getpixel((100, 5)) == (0, 0, 0)


def test_no_decisions_is_no_op(tmp_path: Path) -> None:
    src = _png_solid((100, 100, 100))
    p = tmp_path / "noop.pptx"
    _make_pptx(p, src)
    before = p.read_bytes()
    PptxAdapter().apply_image_redactions(p, {})
    after = p.read_bytes()
    assert before == after
